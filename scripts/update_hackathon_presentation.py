"""
Generate and update the FinExplain 10-slide Hackathon Presentation with exact production data.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path: str):
    prs = Presentation()
    # 16:9 Widescreen slides (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_COLOR = RGBColor(11, 17, 32)        # Deep Navy / Dark Canvas
    CARD_BG = RGBColor(22, 30, 49)        # Elevated Surface Navy
    CARD_BORDER = RGBColor(51, 65, 85)    # Slate Border
    ACCENT_CYAN = RGBColor(6, 182, 212)   # Vibrant Cyan Accent
    ACCENT_EMERALD = RGBColor(16, 185, 129)# Green Success Accent
    ACCENT_AMBER = RGBColor(245, 158, 11) # Warning / Alert Accent
    ACCENT_PURPLE = RGBColor(168, 85, 247)# Purple
    TEXT_WHITE = RGBColor(255, 255, 255)  # Primary White
    TEXT_MUTED = RGBColor(148, 163, 184)  # Secondary Muted Slate
    TEXT_DARK = RGBColor(15, 23, 42)

    blank_layout = prs.slide_layouts[6] # completely blank layout

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="FINEXPLAIN • PRODUCTION ARCHITECTURE"):
        # Top Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_CYAN
        line.line.fill.background()

        # Category Pill / Header
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.48), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(9.5)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.65))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title, body_bullets, accent_color=None, highlight_badge=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = accent_color or CARD_BORDER
        card.line.width = Pt(1.2 if accent_color else 0.8)

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color or TEXT_WHITE
        p_t.space_after = Pt(6)

        for b in body_bullets:
            p_b = tf.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.size = Pt(10)
            p_b.font.color.rgb = TEXT_MUTED
            p_b.space_after = Pt(4)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    add_background(s1)

    # Decorative glow box
    glow = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.1))
    glow.fill.solid()
    glow.fill.fore_color.rgb = CARD_BG
    glow.line.color.rgb = RGBColor(30, 41, 59)
    glow.line.width = Pt(1)

    # Title Text
    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11.0), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "EVIDENCE-FIRST AI FOR COMPLEX CREDIT AGREEMENTS & KFS"
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_CYAN
    p_badge.space_after = Pt(10)

    p_main = tf1.add_paragraph()
    p_main.text = "FinExplain"
    p_main.font.size = Pt(44)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_WHITE
    p_main.space_after = Pt(8)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Zero-Hallucination Financial RAG • Verbatim Page Citations • Deterministic Math"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXT_MUTED

    # Highlights Row on Title Slide
    stat_boxes = [
        ("100% Faithfulness", "Zero hallucinations across 25 benchmark audits", ACCENT_EMERALD),
        ("88.5% Covenant Recall", "Preserves GST, 12-EMI lock-ins & notice terms", ACCENT_CYAN),
        ("90.7% Citation Accuracy", "Deterministic claim-to-chunk page verification", ACCENT_PURPLE),
    ]
    for i, (stitle, ssub, scolor) in enumerate(stat_boxes):
        bx = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2 + i * 3.7), Inches(4.1), Inches(3.5), Inches(1.1))
        bx.fill.solid()
        bx.fill.fore_color.rgb = RGBColor(15, 23, 42)
        bx.line.color.rgb = scolor
        bx.line.width = Pt(1)

        btb = s1.shapes.add_textbox(Inches(1.35 + i * 3.7), Inches(4.2), Inches(3.2), Inches(0.9))
        btf = btb.text_frame
        btf.word_wrap = True
        bp1 = btf.paragraphs[0]
        bp1.text = stitle
        bp1.font.size = Pt(12)
        bp1.font.bold = True
        bp1.font.color.rgb = scolor
        bp2 = btf.add_paragraph()
        bp2.text = ssub
        bp2.font.size = Pt(9.5)
        bp2.font.color.rgb = TEXT_MUTED

    # Team Footer
    foot_tb = s1.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(11.0), Inches(0.6))
    foot_tf = foot_tb.text_frame
    fp = foot_tf.paragraphs[0]
    fp.text = "Team CodeFlex:  Samadhan Mane  •  Sakshi Bhingarkar  •  Pranav Karande   |   Model: Google Gemini 3.5 Flash Lite"
    fp.font.size = Pt(10)
    fp.font.color.rgb = TEXT_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: The Problem
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank_layout)
    add_background(s2)
    add_header(s2, "The Problem: Why Generic RAG Fails on Real Financial Contracts", "01 • PROBLEM STATEMENT")

    add_card(s2, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.2), "1. The 'Lost Qualifier' Flaw", [
        "Generic LLMs summarize text and discard crucial fine print: '+ 18% GST', '12-EMI lock-in', '30-day written notice'.",
        "Result: Borrowers face unexpected penalties of thousands of rupees when prepaying or switching loans.",
        "Traditional RAG benchmark: Drops ~50% of nested contractual conditions."
    ], ACCENT_AMBER)

    add_card(s2, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.2), "2. Arithmetic Hallucinations", [
        "LLMs cannot perform reliable multi-step financial math (EPI schedules, broken period interest, foreclosure charges).",
        "Raw LLM math yields ₹2.25 to ₹50+ errors per calculation.",
        "Fatal for credit underwriters who require exact statutory compliance and balance reconciliation."
    ], RGBColor(239, 68, 68))

    add_card(s2, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2), "3. Opaque Citations & Guessing", [
        "Standard chatbots provide a single generic citation at the end of answers, leaving 50% of claims unverified.",
        "Blindly answers unanswerable queries instead of abstaining when evidence is missing.",
        "Causes fatal legal exposure for regulated lending and audit institutions."
    ], ACCENT_PURPLE)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: Why It Matters (Market & Regulatory Need)
    # ══════════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank_layout)
    add_background(s3)
    add_header(s3, "Why This Matters: Financial Literacy Gaps & Regulatory Scrutiny", "02 • MARKET & REGULATORY URGENCY")

    metrics_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(2.2))
    metrics_card.fill.solid()
    metrics_card.fill.fore_color.rgb = CARD_BG
    metrics_card.line.color.rgb = CARD_BORDER

    tb_m = s3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.8))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    p_mt = tf_m.paragraphs[0]
    p_mt.text = "THE COST OF OPAQUE LOAN AGREEMENTS"
    p_mt.font.size = Pt(12)
    p_mt.font.bold = True
    p_mt.font.color.rgb = ACCENT_CYAN
    p_mt.space_after = Pt(6)

    p_mb = tf_m.add_paragraph()
    p_mb.text = "• 50+ Page Loan Agreements: Borrowers sign multi-tier contracts without discovering hidden lock-ins or reset spreads.\n• RBI Key Fact Statement (KFS) Mandate: Regulators require transparent, standardized APR and fee disclosures across all lenders.\n• 70%+ of Consumer Complaints: Stem from unexpected foreclosure charges, broken-period fees, and interest benchmark variation."
    p_mb.font.size = Pt(10.5)
    p_mb.font.color.rgb = TEXT_MUTED
    p_mb.space_after = Pt(4)

    add_card(s3, Inches(0.8), Inches(4.1), Inches(5.7), Inches(2.8), "Retail & SME Borrowers", [
        "Need instant plain-English explanations with direct page references.",
        "Compare total scenario costs before signing binding legal commitments.",
        "Detect discrepancy between loan sanction letter and master agreement."
    ], ACCENT_EMERALD)

    add_card(s3, Inches(6.8), Inches(4.1), Inches(5.7), Inches(2.8), "Credit Underwriters & Compliance Officers", [
        "Audit Key Fact Statements (KFS) against sanctioned credit limits in seconds.",
        "Automated covenant extraction (DSCR ratios, hypothecation, debt caps).",
        "Deterministic audit trail backed by Human-in-the-Loop review queue."
    ], ACCENT_CYAN)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: The Proposed Solution
    # ══════════════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank_layout)
    add_background(s4)
    add_header(s4, "The Solution: Grounded, Explainable Financial RAG", "03 • PRODUCT SOLUTION")

    steps = [
        ("1. Ingest & Structure", "PyMuPDF sanitizes PDFs, extracts tables into Markdown, and saves core facts to loan_facts.json.", ACCENT_CYAN),
        ("2. Hybrid Retrieve", "Concurrent Pinecone Serverless (384d dense) + Supabase BM25 sparse search for exact terms.", ACCENT_PURPLE),
        ("3. Neural Re-Rank", "ms-marco cross-encoder jointly scores query-chunk pairs, boosting operative clauses (MRR: 0.88).", ACCENT_EMERALD),
        ("4. Deterministic Math", "Python financial math computes APR, EPI, and foreclosure costs with 100% numerical exactness.", ACCENT_AMBER),
        ("5. Verify & Ground", "Claim Verifier checks chunk pages; <30 triggers refusal, <70% routes to HITL compliance queue.", ACCENT_CYAN),
    ]

    for i, (stitle, sbody, scolor) in enumerate(steps):
        add_card(s4, Inches(0.8 + i * 2.38), Inches(1.6), Inches(2.25), Inches(5.2), stitle, [sbody], scolor)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Technical System Architecture
    # ══════════════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank_layout)
    add_background(s5)
    add_header(s5, "Production Architecture: 6-Layer Modular Architecture", "04 • TECHNICAL SYSTEM ARCHITECTURE")

    layers = [
        ("Layer 1: Client Application", "Vite + React 18 + TypeScript + Tailwind CSS | Interactive PDF Viewer & Citation Inspector"),
        ("Layer 2: API Gateway & Security", "FastAPI + SlowAPI Rate Limiting (60 req/min) + InjectionGuard + PiiGuard (PAN/Aadhaar)"),
        ("Layer 3: Parsing & Facts", "PyMuPDF Layout Detection + Markdown Table Extractor + Semantic Chunker (512 tokens)"),
        ("Layer 4: Dual Storage Layer", "Pinecone Serverless Vector DB (Dense Embeddings) + Supabase PostgreSQL (BM25 Sparse Search)"),
        ("Layer 5: Reranking & Context", "Multi-Tier Router (4 Tiers) + ms-marco Cross-Encoder + Clause Context Builder (≤ 2000 tok)"),
        ("Layer 6: Generation & Safety", "Google Gemini 3.5 Flash Lite + Completeness Retry Loop + 7-Dimension Scorer + HITL Queue"),
    ]

    for i, (ltitle, ldesc) in enumerate(layers):
        c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + i * 0.88), Inches(11.733), Inches(0.76))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = ACCENT_CYAN if i >= 4 else CARD_BORDER
        c.line.width = Pt(1.2 if i >= 4 else 0.8)

        tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.65 + i * 0.88), Inches(11.3), Inches(0.65))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        pl1 = tf_l.paragraphs[0]
        pl1.text = ltitle
        pl1.font.size = Pt(11)
        pl1.font.bold = True
        pl1.font.color.rgb = ACCENT_CYAN if i >= 4 else TEXT_WHITE
        pl2 = tf_l.add_paragraph()
        pl2.text = ldesc
        pl2.font.size = Pt(9.5)
        pl2.font.color.rgb = TEXT_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Trust as a Feature (Verification & HITL)
    # ══════════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank_layout)
    add_background(s6)
    add_header(s6, "Trust is a Product Feature: Zero-Hallucination Verification & HITL", "05 • VERIFICATION & GOVERNANCE")

    add_card(s6, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2), "Deterministic Claim Verifier", [
        "Splits LLM generated response into discrete atomic claims.",
        "Asymmetric Word Containment: Measures claim words subset in source chunk (≥ 0.55 threshold).",
        "Coordinate Mapping: Checks that cited [Doc, Page X, Section Y] actually exists on retrieved chunk.",
        "7-Dimension Scorer: Evaluates retrieval similarity, chunk coverage, citation presence, and consistency.",
        "Hard Safety Gate: Automatically refuses delivery if Evidence Score < 30 ('Insufficient Evidence')."
    ], ACCENT_CYAN)

    add_card(s6, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "Human-in-the-Loop (HITL) Governance", [
        "< 70% Confidence Routing Gate: Queries with evidence scores between 30% and 69% route to compliance review queue.",
        "Conflict Detection: Highlights discrepancies between sanction letters (e.g. 10.5%) and master agreements (e.g. 14.0%).",
        "Immutable Audit Logging: Full trace of chunk IDs, re-ranking scores, model prompt, and reviewer approvals.",
        "Regulatory Compliance: Aligned with RBI & financial consumer protection disclosure standards."
    ], ACCENT_EMERALD)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Killer Use Case (Scenario Comparison)
    # ══════════════════════════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(blank_layout)
    add_background(s7)
    add_header(s7, "Killer Scenario: Which Loan is Actually Cheaper?", "06 • SCENARIO SIMULATION ENGINE")

    card_s7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    card_s7.fill.solid()
    card_s7.fill.fore_color.rgb = CARD_BG
    card_s7.line.color.rgb = CARD_BORDER

    # Comparison Table
    rows = [
        ["Evaluation Term", "Product A (Axis LAP)", "Product B (South Indian Bank PL)", "FinExplain Audit Finding"],
        ["Headline Rate", "10.50% p.a. (Fixed)", "11.25% p.a. (Floating spread)", "Product A has lower base rate"],
        ["Processing Fee", "1.00% + 18% GST (₹11,800)", "Fixed ₹2,500 upfront", "Product B is ₹9,300 cheaper upfront"],
        ["Foreclosure Lock-in", "No prepayment for first 12 months", "Prepayment allowed with 30-day notice", "Product A locks borrower for 1 year"],
        ["Foreclosure Charge", "3.00% + GST on prepaid principal", "Zero foreclosure fee after 12 EMIs", "Product B saves ₹35,400 on exit"],
        ["6-Month Prepay Cost", "₹1,24,600 (Total outgo with penalty)", "₹1,14,200 (Total outgo)", "🟢 Product B is ₹10,400 CHEAPER!"]
    ]

    table_shape = s7.shapes.add_table(len(rows), 4, Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.7))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(3.0)
    tbl.columns[2].width = Inches(3.2)
    tbl.columns[3].width = Inches(2.933)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10 if r_idx > 0 else 10.5)
                p.font.bold = (r_idx == 0 or r_idx == len(rows)-1)
                if r_idx == 0:
                    p.font.color.rgb = ACCENT_CYAN
                elif r_idx == len(rows)-1:
                    p.font.color.rgb = ACCENT_EMERALD if c_idx == 3 else TEXT_WHITE
                else:
                    p.font.color.rgb = TEXT_WHITE if c_idx < 3 else TEXT_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Production Evaluation & Benchmark Claims
    # ══════════════════════════════════════════════════════════════════════════
    s8 = prs.slides.add_slide(blank_layout)
    add_background(s8)
    add_header(s8, "Production Evaluation Scorecard: 25-Query Benchmark Results", "07 • BENCHMARK EVALUATION")

    bench_stats = [
        ("100.0%", "Faithfulness Rate", "Zero hallucinations across 25 queries", ACCENT_EMERALD),
        ("88.5%", "Covenant Gen Recall", "Preserves GST, lock-ins & notice", ACCENT_CYAN),
        ("81.7%", "Claim Citation Coverage", "Per-sentence inline citations", ACCENT_PURPLE),
        ("90.7%", "Citation Accuracy", "Verified against chunk page metadata", ACCENT_EMERALD),
        ("4.95s", "P50 Median Latency", "Sub-5s complete verification loop", ACCENT_CYAN),
        ("~$0.000085", "Unit Cost / Query", "Ultra-low Gemini 3.5 Flash Lite economics", ACCENT_EMERALD),
    ]

    for i, (val, title, desc, col) in enumerate(bench_stats):
        col_idx = i % 3
        row_idx = i // 3
        bx = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col_idx * 3.98), Inches(1.6 + row_idx * 2.65), Inches(3.78), Inches(2.45))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_BG
        bx.line.color.rgb = col
        bx.line.width = Pt(1.2)

        tb_b = s8.shapes.add_textbox(Inches(1.0 + col_idx * 3.98), Inches(1.8 + row_idx * 2.65), Inches(3.38), Inches(2.0))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        pb1 = tf_b.paragraphs[0]
        pb1.text = val
        pb1.font.size = Pt(28)
        pb1.font.bold = True
        pb1.font.color.rgb = col
        pb1.space_after = Pt(4)

        pb2 = tf_b.add_paragraph()
        pb2.text = title
        pb2.font.size = Pt(12)
        pb2.font.bold = True
        pb2.font.color.rgb = TEXT_WHITE
        pb2.space_after = Pt(3)

        pb3 = tf_b.add_paragraph()
        pb3.text = desc
        pb3.font.size = Pt(9.5)
        pb3.font.color.rgb = TEXT_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9: Uniqueness & Competitive Advantage
    # ══════════════════════════════════════════════════════════════════════════
    s9 = prs.slides.add_slide(blank_layout)
    add_background(s9)
    add_header(s9, "Competitive Advantage: Why FinExplain is More Than a Document Chatbot", "08 • COMPETITIVE ADVANTAGE")

    comp_rows = [
        ["Capability Dimension", "Generic AI / Chatbots", "Traditional Enterprise RAG", "FinExplain Financial RAG"],
        ["Retrieval Quality", "Single vector search", "Dense vector only (Chroma)", "Hybrid Pinecone + Supabase BM25 (0.88 MRR)"],
        ["Condition Retention", "Summarizes & drops conditions", "Drops 50% of qualifiers", "88.5% Covenant Preservation (Taxes, Lock-ins)"],
        ["Citation Verifiability", "Generic link at bottom", "Chunk-level link", "Per-sentence [Doc, Page X, Section Y] (90.7%)"],
        ["Math & Fee Calculations", "LLM hallucinated numbers", "LLM arithmetic (high error)", "Deterministic Python Financial Engine (₹0.00 MAE)"],
        ["Borderline Uncertainty", "Blindly guesses", "Blindly guesses", "Hard Refusal (<30) + <70% HITL Compliance Queue"],
        ["Security & Privacy", "None (Direct prompt)", "Basic prompt prefix", "SlowAPI Rate Limiter + InjectionGuard + PiiGuard"]
    ]

    tbl_comp_shape = s9.shapes.add_table(len(comp_rows), 4, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tbl_c = tbl_comp_shape.table
    tbl_c.columns[0].width = Inches(2.2)
    tbl_c.columns[1].width = Inches(2.8)
    tbl_c.columns[2].width = Inches(2.8)
    tbl_c.columns[3].width = Inches(3.933)

    for r_idx, row in enumerate(comp_rows):
        for c_idx, val in enumerate(row):
            cell = tbl_c.cell(r_idx, c_idx)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5 if r_idx > 0 else 10.5)
                p.font.bold = (r_idx == 0 or c_idx == 3)
                if r_idx == 0:
                    p.font.color.rgb = ACCENT_CYAN
                elif c_idx == 3:
                    p.font.color.rgb = ACCENT_EMERALD
                else:
                    p.font.color.rgb = TEXT_WHITE if c_idx == 0 else TEXT_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10: Future Roadmap & Engineering Milestones
    # ══════════════════════════════════════════════════════════════════════════
    s10 = prs.slides.add_slide(blank_layout)
    add_background(s10)
    add_header(s10, "Future Roadmap: Expanding the Financial Verification Layer", "09 • FUTURE ROADMAP & SCALABILITY")

    add_card(s10, Inches(0.8), Inches(1.6), Inches(5.7), Inches(2.5), "1. Vision-Based Table Extraction (Q4 2026)", [
        "Integrate multimodal vision parsers (LayoutLMv3 / Gemini Multimodal) for scanned, multi-nested amortization grids.",
        "Zero OCR column-misalignment on complex loan pricing matrices."
    ], ACCENT_CYAN)

    add_card(s10, Inches(6.8), Inches(1.6), Inches(5.7), Inches(2.5), "2. Temporal Contract Diff Graphs", [
        "Automated contract diffing that tracks loan addendums across time.",
        "Master Agreement → Addendum 1 → Sanction Revision with visual covenant change tracking."
    ], ACCENT_PURPLE)

    add_card(s10, Inches(0.8), Inches(4.3), Inches(5.7), Inches(2.5), "3. Real-Time RBI Regulatory Validator", [
        "Automated compliance checkers verifying Key Fact Statements against statutory Reserve Bank of India (RBI) mandates.",
        "Instant alerts on non-compliant penal interest formulas."
    ], ACCENT_EMERALD)

    add_card(s10, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.5), "4. Sub-Second Streaming & Model Fine-Tuning", [
        "Server-Sent Events (SSE) streaming delivering Time-to-First-Token in < 800ms.",
        "Domain-specific fine-tuning on credit contracts pushing Context Recall > 92%."
    ], ACCENT_AMBER)

    prs.save(output_path)
    print(f"[PPTX] Presentation saved successfully at: {output_path}")


if __name__ == "__main__":
    out_target = os.path.abspath(r"d:\Projects\fine-explain\reports\FinExplain_Hackathon_CodeFlex (1).pptx")
    create_presentation(out_target)
